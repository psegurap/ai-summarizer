from pypdf import PdfReader
from docx import Document
from pptx import Presentation


class FileConverter:

    # pdf -----> application/pdf
    @staticmethod
    def convert_PDF(file):
        file_content = ""
        for page in PdfReader(file).pages:
            file_content += page.extract_text(extraction_mode="layout")
        return file_content

    # docx ----> application/vnd.openxmlformats - officedocument.wordprocessingml.document
    @staticmethod
    def convert_DOC(file):
        file_content = ""
        for paragraph in Document(file).paragraphs:
            file_content += f'{paragraph.text}\n'
        return file_content

    # pptx ----> application/vnd.openxmlformats-officedocument.presentationml.presentation
    @staticmethod
    def convert_PPT(file):
        file_content = ""
        for slide in Presentation(file).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    file_content += shape.text
        return file_content.replace('\n', ' ').replace('\r', '')

    # txt ----> text/plain
    @staticmethod
    def convert_TXT(file):
        return file.getvalue().decode("utf-8").replace('\n', ' ').replace('\r', '')

    # md ----> application/octet-stream
    @staticmethod
    def convert_MD(file):
        return file.getvalue().decode("utf-8").replace('\n', ' ').replace('\r', '')
